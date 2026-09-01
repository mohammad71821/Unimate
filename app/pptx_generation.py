import io
import os
import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1A, 0x3C, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_GRAY = RGBColor(0xD8, 0xDE, 0xE8)
BULLET_COLOR = RGBColor(0x2E, 0x6B, 0xB8)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)

FONT_NAME = "Vazirmatn"
_FONTS_DIR = Path(__file__).resolve().parent / "fonts"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MIN_SLIDES = 3
MAX_SLIDES = 20

_ARABIC_TO_PERSIAN = {
    "\u064a": "\u06cc",
    "\u0643": "\u06a9",
    "\u06c0": "\u0647",
}

def _normalize_persian(text: str) -> str:
    if not text:
        return text
    for arabic_ch, persian_ch in _ARABIC_TO_PERSIAN.items():
        text = text.replace(arabic_ch, persian_ch)
    text = re.sub(r" +([.,،؛:!؟])", r"\1", text)
    text = re.sub(r" {2,}", " ", text)
    return text.strip()

def _set_rtl(paragraph) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")

def _style_run(run, size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(f"{{{A_NS}}}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{{{A_NS}}}{tag}")
        el.set("typeface", FONT_NAME)

def _no_line(shape) -> None:
    shape.line.fill.background()

def _add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return box, tf

def _title_font_size(title: str, base: int, long_len: int = 30, longer_len: int = 45) -> int:
    length = len(title or "")
    if length > longer_len:
        return max(base - 10, 18)
    if length > long_len:
        return max(base - 6, 20)
    return base

def _bullet_font_size(bullets: list[str]) -> int:
    total_chars = sum(len(b or "") for b in bullets)
    count = len(bullets)
    if count >= 7 or total_chars > 620:
        return 15
    if count >= 6 or total_chars > 460:
        return 17
    if count >= 5 or total_chars > 320:
        return 19
    if total_chars > 180:
        return 21
    return 23

def _write_bullet(tf, first: bool, text: str, size: int, color: RGBColor) -> None:
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = PP_ALIGN.RIGHT
    _set_rtl(para)
    para.space_after = Pt(12)
    para.line_spacing = 1.25
    text = _normalize_persian(text)
    run_mark = para.add_run()
    run_mark.text = "●  "
    _style_run(run_mark, size, BULLET_COLOR, bold=True)
    run_text = para.add_run()
    run_text.text = text
    _style_run(run_text, size, color)

def build_pptx_bytes(outline: dict) -> bytes:
    slides = outline.get("slides", [])
    if not slides:
        raise ValueError("Outline has no slides")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    first = slides[0]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    title_text = _normalize_persian(first.get("title", ""))
    title_size = _title_font_size(title_text, base=42, long_len=35, longer_len=55)
    _, title_tf = _add_textbox(slide, Inches(1), Inches(2.5), Inches(11.33), Inches(1.8), anchor=MSO_ANCHOR.MIDDLE)
    title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_rtl(p)
    run = p.add_run()
    run.text = title_text
    _style_run(run, title_size, WHITE, bold=True)
    subtitle = _normalize_persian(first.get("subtitle", ""))
    if subtitle:
        _, sub_tf = _add_textbox(slide, Inches(1), Inches(4.4), Inches(11.33), Inches(0.9))
        sub_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = sub_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_rtl(p)
        run = p.add_run()
        run.text = subtitle
        _style_run(run, 18, LIGHT_GRAY)
    accent = slide.shapes.add_shape(1, Inches(11.83), Inches(6.7), Inches(1.2), Inches(0.45))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    _no_line(accent)

    rest = slides[1:] if len(slides) > 1 else slides
    for item in rest:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        header = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.25))
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        _no_line(header)
        accent_line = slide.shapes.add_shape(1, 0, Inches(1.25), SLIDE_W, Inches(0.06))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT
        _no_line(accent_line)
        item_title = _normalize_persian(item.get("title", ""))
        content_title_size = _title_font_size(item_title, base=26, long_len=28, longer_len=42)
        _, title_tf = _add_textbox(slide, Inches(0.6), 0, Inches(12.1), Inches(1.25), anchor=MSO_ANCHOR.MIDDLE)
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = title_tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        _set_rtl(p)
        run = p.add_run()
        run.text = item_title
        _style_run(run, content_title_size, WHITE, bold=True)
        bullets = item.get("bullets", [])
        if bullets:
            bullet_size = _bullet_font_size(bullets)
            _, body_tf = _add_textbox(slide, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.3))
            body_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for i, bullet in enumerate(bullets):
                _write_bullet(body_tf, i == 0, bullet, bullet_size, DARK_GRAY)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
