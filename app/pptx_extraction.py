from pathlib import Path

from pptx import Presentation

PPTX_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

PPTX_EXTENSIONS = {".pptx"}


def is_pptx(content_type: str, filename: str) -> bool:
    if content_type in PPTX_CONTENT_TYPES:
        return True
    return Path(filename).suffix.lower() in PPTX_EXTENSIONS


def extract_text_from_pptx(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    slide_texts = []

    for i, slide in enumerate(presentation.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slide_content = "\n".join(parts)
        block = f"--- اسلاید {i} ---\n{slide_content}"
        if notes:
            block += f"\n[یادداشت اسلاید: {notes}]"
        slide_texts.append(block)

    return "\n\n".join(slide_texts).strip()
